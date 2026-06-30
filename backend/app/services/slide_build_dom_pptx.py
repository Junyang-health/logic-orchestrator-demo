"""
Preview-HTML to editable PPTX export.

This follows the same core idea as tfriedel/claude-office-skills' html2pptx
flow: render the HTML in a browser, read actual DOM geometry and computed
styles, then rebuild the slide as editable PowerPoint objects instead of as a
flat screenshot.
"""

from __future__ import annotations

import subprocess
import tempfile
from pathlib import Path
from typing import Any

from app.services.slide_build_artifacts import (
    extract_slide_inner_from_document,
    read_slide_document,
    render_slide_document,
)
from app.services.slide_build_scene import render_scene_to_pptx_slide, temp_media_dir

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
RENDER_W_PX = 1600
RENDER_H_PX = 900
PREVIEW_TO_PPTX_TIMEOUT_MS = 15_000
DEFAULT_PPT_FONT = "Aptos"


def build_pptx_from_preview_html(session_id: str, framework: dict[str, Any], dest: Path) -> bool:
    """Build an editable PPTX from the exact generated preview HTML.

    Returns False when the environment cannot run browser extraction or no
    preview HTML exists, allowing callers to fall back to pptx_scene/pptx_spec.
    """
    html_docs = _normalized_preview_documents(session_id, framework)
    if not html_docs:
        return False

    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.util import Inches  # type: ignore[import-untyped]
        from playwright.sync_api import sync_playwright  # type: ignore[import-untyped]
    except ImportError:
        return False

    dest.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    chrome = _chrome_executable()
    with tempfile.TemporaryDirectory(prefix="unbox-html2pptx-") as tmp_name:
        tmp_dir = Path(tmp_name)
        html_paths: list[Path] = []
        for idx, doc in enumerate(html_docs, start=1):
            p = tmp_dir / f"slide_{idx:03d}.html"
            p.write_text(doc, encoding="utf-8")
            html_paths.append(p)

        try:
            with sync_playwright() as p:
                browser = _launch_chromium(p.chromium, chrome)
                try:
                    with temp_media_dir() as media_tmp:
                        media_dir = Path(media_tmp)
                        for html_path in html_paths:
                            page = browser.new_page(viewport={"width": RENDER_W_PX, "height": RENDER_H_PX})
                            try:
                                page.goto(html_path.resolve().as_uri(), wait_until="load", timeout=PREVIEW_TO_PPTX_TIMEOUT_MS)
                                page.wait_for_timeout(120)
                                scene = page.evaluate(_EXTRACT_SCENE_JS)
                            finally:
                                page.close()
                            if not isinstance(scene, dict) or not isinstance(scene.get("elements"), list):
                                return False
                            slide = prs.slides.add_slide(blank)
                            render_scene_to_pptx_slide(slide, scene, media_dir)
                finally:
                    browser.close()
        except Exception:
            return False

    prs.save(str(dest))
    return True


def _launch_chromium(chromium: Any, chrome: str | None) -> Any:
    base_kwargs: dict[str, Any] = {
        "headless": True,
        "args": [
            "--disable-gpu",
            "--no-first-run",
            "--no-default-browser-check",
            "--allow-file-access-from-files",
            "--disable-background-networking",
            "--disable-sync",
            "--disable-extensions",
            "--disable-features=Translate,MediaRouter",
        ],
    }
    attempts: list[dict[str, Any]] = []
    if chrome:
        attempts.append({**base_kwargs, "executable_path": chrome})
    attempts.append(base_kwargs)

    last_err: Exception | None = None
    for kwargs in attempts:
        try:
            return chromium.launch(**kwargs)
        except Exception as exc:
            last_err = exc
    if last_err:
        raise last_err
    raise RuntimeError("Chromium launch failed")


def _framework_slide_ids(framework: dict[str, Any]) -> list[str]:
    raw = framework.get("slides")
    if not isinstance(raw, list):
        return []
    out: list[str] = []
    for idx, slide in enumerate(raw, start=1):
        if isinstance(slide, dict):
            sid = str(slide.get("id") or "").strip()
            out.append(sid or f"slide_{idx}")
    return out


def _normalized_preview_documents(session_id: str, framework: dict[str, Any]) -> list[str]:
    docs: list[str] = []
    for sid in _framework_slide_ids(framework):
        full = read_slide_document(session_id, sid)
        if not full:
            continue
        inner = extract_slide_inner_from_document(full)
        docs.append(render_slide_document(inner) if inner else full)
    return docs


def _chrome_executable() -> str | None:
    candidates = [
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
        "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
        "google-chrome",
        "chromium",
        "chromium-browser",
    ]
    for cand in candidates:
        if "/" in cand:
            if Path(cand).is_file():
                return cand
            continue
        try:
            subprocess.run(
                [cand, "--version"],
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=5,
            )
            return cand
        except Exception:
            continue
    return None


_EXTRACT_SCENE_JS = r"""
() => {
  const slide = document.querySelector('.slide') || document.body;
  const slideRect = slide.getBoundingClientRect();
  const sw = Math.max(slideRect.width, 1);
  const sh = Math.max(slideRect.height, 1);
  const sx = 1600 / sw;
  const sy = 900 / sh;
  const elements = [];

  function parseColor(css) {
    if (!css || css === 'transparent') return null;
    const m = String(css).match(/rgba?\(([^)]+)\)/i);
    if (!m) return String(css);
    const parts = m[1].split(',').map(x => x.trim());
    const r = Number(parts[0] || 0), g = Number(parts[1] || 0), b = Number(parts[2] || 0);
    const a = parts.length >= 4 ? Number(parts[3]) : 1;
    return { r, g, b, a: Number.isFinite(a) ? a : 1 };
  }

  function hexFromRgb(c) {
    const h = n => Math.max(0, Math.min(255, Math.round(n))).toString(16).padStart(2, '0');
    return `#${h(c.r)}${h(c.g)}${h(c.b)}`;
  }

  const slideStyle = getComputedStyle(slide);
  const slideBg = parseColor(slideStyle.backgroundColor) || { r: 255, g: 255, b: 255, a: 1 };
  const baseBg = slideBg.a < 0.05 ? { r: 255, g: 255, b: 255, a: 1 } : slideBg;

  function color(css, bg = baseBg) {
    const c = parseColor(css);
    if (!c || c.a < 0.03) return null;
    if (c.a >= 0.99) return hexFromRgb(c);
    return hexFromRgb({
      r: c.r * c.a + bg.r * (1 - c.a),
      g: c.g * c.a + bg.g * (1 - c.a),
      b: c.b * c.a + bg.b * (1 - c.a),
    });
  }

  function bounds(rect) {
    return {
      x: (rect.left - slideRect.left) * sx,
      y: (rect.top - slideRect.top) * sy,
      w: rect.width * sx,
      h: rect.height * sy,
    };
  }

  function visible(el, rect, cs) {
    if (!rect || rect.width < 1 || rect.height < 1) return false;
    if (cs.display === 'none' || cs.visibility === 'hidden') return false;
    if (Number(cs.opacity || 1) < 0.04) return false;
    return true;
  }

  function directText(el) {
    const direct = Array.from(el.childNodes)
      .filter(n => n.nodeType === Node.TEXT_NODE)
      .map(n => n.textContent || '')
      .join(' ')
      .replace(/\s+/g, ' ')
      .trim();
    if (direct) return direct;
    const hasTextChild = Array.from(el.children).some(c => (c.innerText || '').trim());
    return hasTextChild ? '' : (el.innerText || '').replace(/\s+/g, ' ').trim();
  }

  function weight(cs) {
    const fw = String(cs.fontWeight || '');
    return fw === 'bold' || Number(fw) >= 600;
  }

  function pptFontFamily(cs) {
    const raw = (cs.fontFamily || '').split(',')[0].replace(/["']/g, '').trim();
    if (!raw || raw === 'system-ui' || raw === '-apple-system' || raw.toLowerCase().includes('segoe ui')) {
      return '""" + DEFAULT_PPT_FONT + r"""';
    }
    return raw;
  }

  function cssPxToPt(px) {
    const renderedPxPerIn = Math.max(sw / 13.333, sh / 7.5, 1);
    return px * 72 / renderedPxPerIn;
  }

  function addElementFor(el) {
    if (el === slide) return;
    const tag = el.tagName.toLowerCase();
    const rect = el.getBoundingClientRect();
    const cs = getComputedStyle(el);
    if (!visible(el, rect, cs)) return;
    const b = bounds(rect);
    if (b.x > 1602 || b.y > 902 || b.x + b.w < -2 || b.y + b.h < -2) return;

    const fill = color(cs.backgroundColor);
    const borderColor = color(cs.borderTopColor);
    const borderWidth = Math.max(
      parseFloat(cs.borderTopWidth || 0),
      parseFloat(cs.borderRightWidth || 0),
      parseFloat(cs.borderBottomWidth || 0),
      parseFloat(cs.borderLeftWidth || 0)
    );
    if (fill || (borderColor && borderWidth >= 0.5)) {
      const radius = Math.max(
        parseFloat(cs.borderTopLeftRadius || 0),
        parseFloat(cs.borderTopRightRadius || 0),
        parseFloat(cs.borderBottomRightRadius || 0),
        parseFloat(cs.borderBottomLeftRadius || 0)
      ) * ((sx + sy) / 2);
      elements.push({
        type: radius > 2 ? 'round_rect' : 'rect',
        ...b,
        fill: fill || 'transparent',
        stroke: borderColor || 'transparent',
        strokeWidth: borderColor ? Math.max(0.25, borderWidth * ((sx + sy) / 2)) : 0,
        radius,
      });
    }

    if (tag === 'img') {
      const src = el.currentSrc || el.src || '';
      if (src.startsWith('data:image/') || src.startsWith('file://')) {
        elements.push({ type: 'image', ...b, src, fit: cs.objectFit || 'contain' });
      }
      return;
    }

    if (tag === 'svg') {
      elements.push({ type: 'svg', ...b, svg: el.outerHTML });
      return;
    }

    let txt = directText(el);
    if (tag === 'li' && txt && cs.listStyleType && cs.listStyleType !== 'none') txt = `• ${txt}`;
    if (txt) {
      const cssFontPx = parseFloat(cs.fontSize || '16');
      const fontSizePt = cssPxToPt(cssFontPx);
      const cssLinePx = cs.lineHeight === 'normal' ? cssFontPx * 1.2 : parseFloat(cs.lineHeight || `${cssFontPx * 1.2}`);
      const lineHeightPt = cssPxToPt(cssLinePx);
      elements.push({
        type: 'text',
        ...b,
        text: txt,
        fontSize: fontSizePt,
        lineHeight: lineHeightPt,
        fontFace: pptFontFamily(cs),
        color: color(cs.color) || '#111827',
        bold: weight(cs),
        italic: cs.fontStyle === 'italic',
        align: cs.textAlign || 'left',
        valign: cs.alignItems === 'center' ? 'middle' : 'top',
      });
    }
  }

  Array.from(slide.querySelectorAll('*')).forEach(addElementFor);
  return {
    canvas: { width: 1600, height: 900 },
    background: color(slideStyle.backgroundColor) || '#ffffff',
    elements: elements.slice(0, 320),
  };
}
"""
