"""
Render editable PPTX using per-slide ``pptx_spec`` (from slide generation) merged with framework rows.

Falls back to framework title / subtitle / main when no spec is stored.
"""

from __future__ import annotations

import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from app.services.slide_build_artifacts import extract_slide_inner_from_document, read_slide_document, read_slide_manifest
from app.services.slide_build_scene import has_scene, render_scene_to_pptx_slide, temp_media_dir
from app.services.slide_build_session_prefs import load_build_preferences
from app.services.document_style import get_document_style, normalize_deck_style

LAYOUTS = frozenset(
    {
        "title_body",
        "title_only",
        "title_subtitle_body",
        "two_column",
        "content_table",
        "kpi_strip",
        "comparison_board",
        "process_flow",
        "timeline",
        "matrix_2x2",
        "status_board",
    }
)


@dataclass(frozen=True)
class PptxSlidePalette:
    slide_bg: tuple[int, int, int]
    title_rgb: tuple[int, int, int]
    subtitle_rgb: tuple[int, int, int]
    body_rgb: tuple[int, int, int]
    accent_rgb: tuple[int, int, int]
    muted_rgb: tuple[int, int, int]
    card_bg: tuple[int, int, int]
    card_border_rgb: tuple[int, int, int]
    success_rgb: tuple[int, int, int]
    danger_rgb: tuple[int, int, int]
    title_pt: int
    subtitle_pt: int
    body_pt: int


@dataclass(frozen=True)
class PreviewComparisonPanel:
    coverage_rows: tuple[tuple[str, str], ...]
    insight_rows: tuple[str, ...]
    takeaway: str


def _normalize_deck_style(value: str | None) -> str:
    return normalize_deck_style(value)


def _slides_list(framework: dict[str, Any]) -> list[dict[str, Any]]:
    raw = framework.get("slides")
    if not isinstance(raw, list):
        return []
    out: list[dict[str, Any]] = []
    for s in raw:
        if isinstance(s, dict):
            out.append(s)
    return out


def _bullet_items_from_spec(spec: dict[str, Any]) -> list[tuple[str, int]]:
    raw = spec.get("body_bullets")
    if not isinstance(raw, list):
        return []
    out: list[tuple[str, int]] = []
    for item in raw:
        if isinstance(item, str) and item.strip():
            out.append((item.strip(), 0))
            continue
        if isinstance(item, dict):
            txt = str(item.get("text") or "").strip()
            if not txt:
                continue
            try:
                lvl = int(item.get("level", 0))
            except (TypeError, ValueError):
                lvl = 0
            out.append((txt, max(0, min(lvl, 3))))
    return out


def _right_column_bullets(spec: dict[str, Any]) -> list[tuple[str, int]]:
    raw = spec.get("column_right_bullets")
    if not isinstance(raw, list):
        return []
    out: list[tuple[str, int]] = []
    for item in raw:
        if isinstance(item, str) and item.strip():
            out.append((item.strip(), 0))
            continue
        if isinstance(item, dict):
            txt = str(item.get("text") or "").strip()
            if not txt:
                continue
            try:
                lvl = int(item.get("level", 0))
            except (TypeError, ValueError):
                lvl = 0
            out.append((txt, max(0, min(lvl, 3))))
    return out


def _main_lines_from_framework(main: str) -> list[tuple[str, int]]:
    lines = [ln.strip() for ln in main.split("\n") if ln.strip()]
    return [(ln, 0) for ln in lines]


def _coerce_spec(raw: Any) -> dict[str, Any] | None:
    if raw is None:
        return None
    if isinstance(raw, dict):
        return raw
    if isinstance(raw, str) and raw.strip():
        try:
            parsed = json.loads(raw)
            return parsed if isinstance(parsed, dict) else None
        except json.JSONDecodeError:
            return None
    return None


def _load_meta(session_id: str | None, slide_id: str) -> dict[str, Any]:
    if not session_id or not slide_id.strip():
        return {}
    meta = read_slide_manifest(session_id, slide_id)
    return dict(meta) if meta else {}


def _load_preview_inner(session_id: str | None, slide_id: str) -> str:
    if not session_id or not slide_id.strip():
        return ""
    full = read_slide_document(session_id, slide_id)
    if not full:
        return ""
    return extract_slide_inner_from_document(full)


def _extract_preview_comparison_panel(session_id: str | None, slide_id: str) -> PreviewComparisonPanel | None:
    inner = _load_preview_inner(session_id, slide_id)
    if not inner or "能力覆盖度" not in inner or "独家能力示意" not in inner:
        return None

    coverage_matches = re.findall(
        r"<span>([^<]+)</span>\s*<span>([^<]+)</span>",
        inner,
        flags=re.IGNORECASE,
    )
    coverage_rows: list[tuple[str, str]] = []
    seen_labels: set[str] = set()
    for label, value in coverage_matches:
        clean_label = re.sub(r"\s+", " ", label).strip()
        clean_value = re.sub(r"\s+", " ", value).strip()
        if not clean_label or not clean_value:
            continue
        if clean_label in {"NotebookLM", "其他竞品", "OpenTactics"} and clean_label not in seen_labels:
            coverage_rows.append((clean_label, clean_value))
            seen_labels.add(clean_label)

    insight_rows: list[str] = []
    for item in re.findall(r"border-radius:50%;background:[^;]+;'\s*></span>([^<]+)</div>", inner, flags=re.IGNORECASE):
        clean = re.sub(r"\s+", " ", item).strip()
        if clean:
            insight_rows.append(clean)

    takeaway = ""
    m_takeaway = re.search(
        r"text-align:center;font-size:14px;font-weight:600;color:#166534;'>([^<]+)</div>",
        inner,
        flags=re.IGNORECASE,
    )
    if m_takeaway:
        takeaway = re.sub(r"\s+", " ", m_takeaway.group(1)).strip()

    if not coverage_rows and not insight_rows and not takeaway:
        return None
    return PreviewComparisonPanel(
        coverage_rows=tuple(coverage_rows[:3]),
        insight_rows=tuple(insight_rows[:6]),
        takeaway=takeaway[:240],
    )


def _merged_slide_view(
    slide_fw: dict[str, Any],
    spec: dict[str, Any] | None,
    framework_deck_style: str,
) -> dict[str, Any]:
    s = spec or {}
    title = str(s.get("title") or slide_fw.get("title") or "Slide").strip() or "Slide"
    subtitle = str(s.get("subtitle") or slide_fw.get("subtitle") or "").strip()
    takeaway = str(s.get("takeaway") or "").strip()
    footer = str(s.get("footer_source") or "").strip()
    layout = str(s.get("layout") or "title_body").strip().lower()
    if layout not in LAYOUTS:
        layout = "title_body"
    deck_style = _normalize_deck_style(str(s.get("deck_style") or framework_deck_style))
    bullets = _bullet_items_from_spec(s)
    if not bullets:
        bullets = _main_lines_from_framework(str(slide_fw.get("main") or ""))
    visual = s.get("visual") if isinstance(s.get("visual"), dict) else {}
    metrics = s.get("metrics") if isinstance(s.get("metrics"), list) else []
    right_col = _right_column_bullets(s)
    return {
        "title": title,
        "subtitle": subtitle,
        "takeaway": takeaway,
        "footer_source": footer,
        "layout": layout,
        "deck_style": deck_style,
        "bullets": bullets,
        "right_col": right_col,
        "visual": visual,
        "metrics": metrics,
    }


def merged_slide_views(session_id: str | None, framework: dict[str, Any]) -> list[dict[str, Any]]:
    """Aligned list of merged views (order = framework slides) for PPTX/PDF."""
    fw_style = _normalize_deck_style(str(framework.get("deck_style")))
    slides = _slides_list(framework)
    if not slides:
        slides = [{"id": "__empty__", "title": "Empty deck", "subtitle": "", "main": "Add slides to the framework."}]
    views: list[dict[str, Any]] = []
    for fw_slide in slides:
        sid = str(fw_slide.get("id") or "")
        meta = _load_meta(session_id, sid)
        raw_spec = meta.get("pptx_spec") if isinstance(meta, dict) else None
        spec = _coerce_spec(raw_spec)
        views.append(_merged_slide_view(fw_slide, spec, fw_style))
    return views


def _apply_background(slide: Any, rgb: tuple[int, int, int]) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])


def _fill_bullets(tf: Any, items: list[tuple[str, int]], *, font_pt: int, color: Any) -> None:
    from pptx.util import Pt  # type: ignore[import-untyped]

    tf.clear()
    if not items:
        tf.paragraphs[0].text = ""
        return
    first = True
    for text, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text[:4000]
        p.level = level
        p.font.size = Pt(font_pt)
        p.font.color.rgb = color


def _layout_by_name(prs: Any, substrs: tuple[str, ...]) -> Any:
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if all(s in name for s in substrs):
            return layout
    return None


def _surface_variant(session_id: str | None) -> str:
    if not session_id:
        return "light"
    prefs = load_build_preferences(session_id)
    design = prefs.get("design") if isinstance(prefs.get("design"), dict) else {}
    surface = str(design.get("surface") or "").strip().lower()
    return surface if surface in {"light", "dark", "glass", "mono"} else "light"


def _palette_for(style: str, surface: str = "light") -> PptxSlidePalette:
    shared = get_document_style(style, surface)
    return PptxSlidePalette(
        slide_bg=shared.slide_bg,
        title_rgb=shared.title_rgb,
        subtitle_rgb=shared.subtitle_rgb,
        body_rgb=shared.body_rgb,
        accent_rgb=shared.accent_rgb,
        muted_rgb=shared.muted_rgb,
        card_bg=shared.card_bg,
        card_border_rgb=shared.card_border_rgb,
        success_rgb=shared.success_rgb,
        danger_rgb=shared.danger_rgb,
        title_pt=shared.title_pt,
        subtitle_pt=shared.subtitle_pt,
        body_pt=shared.body_pt,
    )


def _add_table_to_slide(
    slide: Any,
    headers: list[str],
    rows: list[list[str]],
    palette: Any,
    *,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    highlight_ot: bool = False,
) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.util import Inches, Pt  # type: ignore[import-untyped]

    if not headers:
        return
    cols = len(headers)
    body_rows = [r for r in rows if isinstance(r, list)][:12]
    nrows = 1 + len(body_rows)
    graphic = slide.shapes.add_table(nrows, cols, left, top, width, height)
    tbl = graphic.table
    hdr_fill = RGBColor(palette.accent_rgb[0], palette.accent_rgb[1], palette.accent_rgb[2])
    hdr_neutral_fill = RGBColor(241, 245, 249) if palette.slide_bg[0] > 200 else RGBColor(*palette.card_bg)
    fg = RGBColor(255, 255, 255)
    fg_dark = RGBColor(51, 65, 85)
    fg_body = RGBColor(palette.body_rgb[0], palette.body_rgb[1], palette.body_rgb[2])
    success = RGBColor(palette.success_rgb[0], palette.success_rgb[1], palette.success_rgb[2])
    danger = RGBColor(palette.danger_rgb[0], palette.danger_rgb[1], palette.danger_rgb[2])
    accent = RGBColor(palette.accent_rgb[0], palette.accent_rgb[1], palette.accent_rgb[2])
    card_fill = RGBColor(palette.card_bg[0], palette.card_bg[1], palette.card_bg[2])

    for c in range(cols):
        tbl.columns[c].width = int(width / cols)

    for c, h in enumerate(headers[:cols]):
        cell = tbl.cell(0, c)
        cell.text = str(h)[:200]
        cell.fill.solid()
        if highlight_ot and c == cols - 1:
            cell.fill.fore_color.rgb = success
        else:
            cell.fill.fore_color.rgb = hdr_neutral_fill if highlight_ot else hdr_fill
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(palette.body_pt - 1)
        p.font.color.rgb = fg if (not highlight_ot or c == cols - 1) else fg_dark

    for r, row in enumerate(body_rows, start=1):
        for c in range(cols):
            val = ""
            if c < len(row):
                val = str(row[c])
            cell = tbl.cell(r, c)
            cell.text = val[:600]
            cell.fill.solid()
            cell.fill.fore_color.rgb = card_fill
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(palette.body_pt - 1)
            low = val.strip().lower()
            positive = (
                "✅" in val
                or "✔" in val
                or any(term in val for term in ("独家", "原子级", "标杆级", "专属", "内置", "变量调节器", "逻辑资产流转", "One-Click", "API+HID"))
            )
            negative = (
                "✕" in val
                or "✖" in val
                or "×" in val
                or low.startswith("no ")
                or any(term in val for term in ("不支持", "有限", "弱", "混合/全局", "共享笔记", "人工", "无法定量", "需配置"))
                or low in {"无", "pdf/web"}
            )
            if positive:
                p.font.color.rgb = success
                p.font.bold = True
            elif negative:
                p.font.color.rgb = danger
                p.font.bold = True
            elif c == cols - 1:
                p.font.color.rgb = success if highlight_ot else accent
                p.font.bold = True
            else:
                p.font.color.rgb = fg_body

    for ridx in range(nrows):
        for cidx in range(cols):
            try:
                cell = tbl.cell(ridx, cidx)
                cell.margin_left = Inches(0.08)
                cell.margin_right = Inches(0.08)
                cell.margin_top = Inches(0.04)
                cell.margin_bottom = Inches(0.04)
            except AttributeError:
                pass


def _add_card(
    slide: Any,
    *,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    palette: PptxSlidePalette,
    radius: float = 0.16,
) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore[import-untyped]
    from pptx.util import Inches  # type: ignore[import-untyped]

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*palette.card_bg)
    shape.line.color.rgb = RGBColor(*palette.card_border_rgb)
    shape.line.width = Inches(0.015)


def _add_textbox(
    slide: Any,
    *,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    text: str,
    font_pt: int,
    rgb: tuple[int, int, int],
    bold: bool = False,
) -> Any:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.util import Pt  # type: ignore[import-untyped]

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_pt)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*rgb)
    return box


def _fit_font_size(text: str, base_pt: int, *, thresholds: list[tuple[int, int]]) -> int:
    size = base_pt
    for limit, candidate in thresholds:
        if len(text) > limit:
            size = min(size, candidate)
    return max(9, size)


def _add_bullet_box(
    slide: Any,
    *,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    items: list[tuple[str, int]],
    palette: PptxSlidePalette,
) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.util import Inches, Pt  # type: ignore[import-untyped]

    _add_card(slide, left=left, top=top, width=width, height=height, palette=palette)
    box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.14), width - Inches(0.36), height - Inches(0.28))
    tf = box.text_frame
    tf.clear()
    first = True
    accent = RGBColor(*palette.accent_rgb)
    body = RGBColor(*palette.body_rgb)
    for text, level in items[:8]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        clean = text.strip()
        prefix = ""
        suffix = clean
        if ":" in clean and len(clean.split(":", 1)[0]) <= 28:
            prefix, suffix = clean.split(":", 1)
            suffix = suffix.strip()
            prefix = prefix.strip() + ":"
        if prefix:
            r1 = p.add_run()
            r1.text = prefix + " "
            r1.font.bold = True
            r1.font.size = Pt(palette.body_pt)
            r1.font.color.rgb = accent
            r2 = p.add_run()
            r2.text = suffix[:800]
            r2.font.size = Pt(palette.body_pt)
            r2.font.color.rgb = body
        else:
            p.text = clean[:900]
            p.font.size = Pt(palette.body_pt)
            p.font.color.rgb = body
        if level > 0:
            p.level = min(level, 2)


def _add_comparison_board(
    slide: Any,
    *,
    vis: dict[str, Any],
    bullets: list[tuple[str, int]],
    palette: PptxSlidePalette,
) -> None:
    from pptx.util import Inches  # type: ignore[import-untyped]

    columns = vis.get("columns") if isinstance(vis.get("columns"), list) else []
    parsed: list[dict[str, Any]] = [c for c in columns[:3] if isinstance(c, dict)]
    if len(parsed) < 2:
        _add_bullet_box(
            slide,
            left=Inches(0.72),
            top=Inches(1.88),
            width=Inches(11.9),
            height=Inches(4.6),
            items=bullets or [("Add comparison columns to pptx_spec.visual.columns.", 0)],
            palette=palette,
        )
        return

    if bullets:
        _add_bullet_box(
            slide,
            left=Inches(0.72),
            top=Inches(1.88),
            width=Inches(3.15),
            height=Inches(4.55),
            items=bullets[:6],
            palette=palette,
        )
        board_left = Inches(4.05)
        board_width = Inches(8.57)
    else:
        board_left = Inches(0.72)
        board_width = Inches(11.9)

    gap = Inches(0.24)
    col_w = (board_width - gap * (len(parsed) - 1)) / len(parsed)
    for idx, col in enumerate(parsed):
        left = board_left + idx * (col_w + gap)
        _add_card(slide, left=left, top=Inches(1.88), width=col_w, height=Inches(4.55), palette=palette)
        _add_textbox(
            slide,
            left=left + Inches(0.18),
            top=Inches(2.04),
            width=col_w - Inches(0.36),
            height=Inches(0.28),
            text=str(col.get("title") or f"Column {idx + 1}")[:60],
            font_pt=palette.body_pt + 1,
            rgb=palette.accent_rgb,
            bold=True,
        )
        sub = str(col.get("subtitle") or "").strip()
        if sub:
            _add_textbox(
                slide,
                left=left + Inches(0.18),
                top=Inches(2.34),
                width=col_w - Inches(0.36),
                height=Inches(0.24),
                text=sub[:90],
                font_pt=max(10, palette.body_pt - 2),
                rgb=palette.subtitle_rgb,
            )
        items_raw = col.get("items") if isinstance(col.get("items"), list) else []
        items: list[tuple[str, int]] = []
        for item in items_raw[:6]:
            if isinstance(item, str) and item.strip():
                items.append((item.strip(), 0))
        _add_bullet_box(
            slide,
            left=left + Inches(0.12),
            top=Inches(2.72),
            width=col_w - Inches(0.24),
            height=Inches(3.45),
            items=items or [("Add column points", 0)],
            palette=palette,
        )


def _add_process_flow(
    slide: Any,
    *,
    vis: dict[str, Any],
    bullets: list[tuple[str, int]],
    palette: PptxSlidePalette,
) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR  # type: ignore[import-untyped]
    from pptx.util import Inches, Pt  # type: ignore[import-untyped]

    steps_raw = vis.get("steps") if isinstance(vis.get("steps"), list) else []
    steps = [s for s in steps_raw[:5] if isinstance(s, dict)]
    if not steps:
        _add_bullet_box(
            slide,
            left=Inches(0.72),
            top=Inches(1.88),
            width=Inches(11.9),
            height=Inches(4.6),
            items=bullets or [("Add process steps to pptx_spec.visual.steps.", 0)],
            palette=palette,
        )
        return

    if bullets:
        _add_bullet_box(
            slide,
            left=Inches(0.72),
            top=Inches(1.88),
            width=Inches(3.45),
            height=Inches(4.6),
            items=bullets[:6],
            palette=palette,
        )
        flow_left = Inches(4.42)
        flow_width = Inches(8.2)
    else:
        flow_left = Inches(0.72)
        flow_width = Inches(11.9)

    gap = Inches(0.18)
    step_w = (flow_width - gap * (len(steps) - 1)) / len(steps)
    top = Inches(2.34)
    card_h = Inches(2.6)
    for idx, step in enumerate(steps):
        left = flow_left + idx * (step_w + gap)
        _add_card(slide, left=left, top=top, width=step_w, height=card_h, palette=palette)
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.12), top + Inches(0.12), Inches(0.34), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(*palette.accent_rgb)
        badge.line.color.rgb = RGBColor(*palette.accent_rgb)
        tf = badge.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = str(idx + 1)
        p.font.size = Pt(max(10, palette.body_pt - 1))
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        _add_textbox(
            slide,
            left=left + Inches(0.54),
            top=top + Inches(0.12),
            width=step_w - Inches(0.68),
            height=Inches(0.32),
            text=str(step.get("title") or f"Step {idx + 1}")[:50],
            font_pt=palette.body_pt,
            rgb=palette.title_rgb,
            bold=True,
        )
        detail = str(step.get("detail") or step.get("subtitle") or "").strip()
        if detail:
            _add_textbox(
                slide,
                left=left + Inches(0.16),
                top=top + Inches(0.58),
                width=step_w - Inches(0.32),
                height=Inches(1.3),
                text=detail[:180],
                font_pt=max(10, palette.body_pt - 2),
                rgb=palette.body_rgb,
            )
        outcome = str(step.get("outcome") or "").strip()
        if outcome:
            _add_textbox(
                slide,
                left=left + Inches(0.16),
                top=top + Inches(1.92),
                width=step_w - Inches(0.32),
                height=Inches(0.34),
                text=outcome[:80],
                font_pt=max(10, palette.body_pt - 2),
                rgb=palette.accent_rgb,
                bold=True,
            )
        if idx < len(steps) - 1:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                left + step_w,
                top + card_h / 2,
                left + step_w + gap,
                top + card_h / 2,
            )
            conn.line.color.rgb = RGBColor(*palette.accent_rgb)
            conn.line.width = Inches(0.02)


def _add_matrix_2x2(
    slide: Any,
    *,
    vis: dict[str, Any],
    bullets: list[tuple[str, int]],
    palette: PptxSlidePalette,
) -> None:
    from pptx.util import Inches  # type: ignore[import-untyped]

    matrix = vis.get("matrix") if isinstance(vis.get("matrix"), dict) else {}
    quadrants = matrix.get("quadrants") if isinstance(matrix.get("quadrants"), list) else []
    cells = [c for c in quadrants[:4] if isinstance(c, dict)]
    if len(cells) < 4:
        cells = [
            {"title": "Scale", "detail": "High impact / high confidence"},
            {"title": "Test", "detail": "High impact / low confidence"},
            {"title": "Maintain", "detail": "Lower impact / high confidence"},
            {"title": "Watch", "detail": "Lower impact / low confidence"},
        ]

    if bullets:
        _add_bullet_box(
            slide,
            left=Inches(0.72),
            top=Inches(1.88),
            width=Inches(3.15),
            height=Inches(4.55),
            items=bullets[:6],
            palette=palette,
        )
        matrix_left = Inches(4.16)
        matrix_width = Inches(8.46)
    else:
        matrix_left = Inches(1.0)
        matrix_width = Inches(11.0)
    matrix_top = Inches(2.0)
    matrix_height = Inches(4.25)
    _add_card(slide, left=matrix_left, top=matrix_top, width=matrix_width, height=matrix_height, palette=palette)

    h_gap = Inches(0.18)
    v_gap = Inches(0.18)
    cell_w = (matrix_width - Inches(0.46) - h_gap) / 2
    cell_h = (matrix_height - Inches(0.64) - v_gap) / 2
    start_left = matrix_left + Inches(0.22)
    start_top = matrix_top + Inches(0.22)
    for idx, cell in enumerate(cells[:4]):
        row = idx // 2
        col = idx % 2
        left = start_left + col * (cell_w + h_gap)
        top = start_top + row * (cell_h + v_gap)
        _add_card(slide, left=left, top=top, width=cell_w, height=cell_h, palette=palette)
        _add_textbox(
            slide,
            left=left + Inches(0.14),
            top=top + Inches(0.12),
            width=cell_w - Inches(0.28),
            height=Inches(0.28),
            text=str(cell.get("title") or f"Quadrant {idx + 1}")[:50],
            font_pt=palette.body_pt,
            rgb=palette.accent_rgb,
            bold=True,
        )
        _add_textbox(
            slide,
            left=left + Inches(0.14),
            top=top + Inches(0.48),
            width=cell_w - Inches(0.28),
            height=cell_h - Inches(0.62),
            text=str(cell.get("detail") or cell.get("summary") or "")[:180],
            font_pt=max(10, palette.body_pt - 2),
            rgb=palette.body_rgb,
        )

    axis_x = str(matrix.get("x_axis") or "Lower → Higher").strip()
    axis_y = str(matrix.get("y_axis") or "Lower → Higher").strip()
    _add_textbox(
        slide,
        left=matrix_left + Inches(2.8),
        top=matrix_top + matrix_height + Inches(0.02),
        width=Inches(3.2),
        height=Inches(0.22),
        text=axis_x[:80],
        font_pt=max(10, palette.body_pt - 3),
        rgb=palette.subtitle_rgb,
        bold=True,
    )
    _add_textbox(
        slide,
        left=matrix_left - Inches(0.1),
        top=matrix_top + Inches(1.8),
        width=Inches(0.9),
        height=Inches(0.24),
        text=axis_y[:80],
        font_pt=max(10, palette.body_pt - 3),
        rgb=palette.subtitle_rgb,
        bold=True,
    )


def _add_timeline(
    slide: Any,
    *,
    vis: dict[str, Any],
    bullets: list[tuple[str, int]],
    palette: PptxSlidePalette,
) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR  # type: ignore[import-untyped]
    from pptx.util import Inches, Pt  # type: ignore[import-untyped]

    events_raw = vis.get("events") if isinstance(vis.get("events"), list) else []
    events = [e for e in events_raw[:5] if isinstance(e, dict)]
    if not events:
        _add_bullet_box(
            slide,
            left=Inches(0.72),
            top=Inches(1.88),
            width=Inches(11.9),
            height=Inches(4.6),
            items=bullets or [("Add timeline events to pptx_spec.visual.events.", 0)],
            palette=palette,
        )
        return

    if bullets:
        _add_bullet_box(
            slide,
            left=Inches(0.72),
            top=Inches(1.88),
            width=Inches(3.15),
            height=Inches(4.55),
            items=bullets[:6],
            palette=palette,
        )
        lane_left = Inches(4.14)
        lane_width = Inches(8.48)
    else:
        lane_left = Inches(0.72)
        lane_width = Inches(11.9)
    lane_top = Inches(4.0)
    lane_mid_y = lane_top + Inches(0.04)
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        lane_left,
        lane_mid_y,
        lane_left + lane_width,
        lane_mid_y,
    )
    conn.line.color.rgb = RGBColor(*palette.accent_rgb)
    conn.line.width = Inches(0.025)

    gap = lane_width / max(len(events), 1)
    for idx, event in enumerate(events):
        center_x = lane_left + gap * idx + gap / 2
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, center_x - Inches(0.08), lane_mid_y - Inches(0.08), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(*palette.accent_rgb)
        dot.line.color.rgb = RGBColor(*palette.accent_rgb)
        card_top = Inches(2.2) if idx % 2 == 0 else Inches(4.22)
        _add_card(
            slide,
            left=center_x - Inches(0.92),
            top=card_top,
            width=Inches(1.84),
            height=Inches(1.32),
            palette=palette,
        )
        _add_textbox(
            slide,
            left=center_x - Inches(0.82),
            top=card_top + Inches(0.1),
            width=Inches(1.64),
            height=Inches(0.22),
            text=str(event.get("label") or event.get("date") or f"Event {idx + 1}")[:40],
            font_pt=max(10, palette.body_pt - 2),
            rgb=palette.accent_rgb,
            bold=True,
        )
        _add_textbox(
            slide,
            left=center_x - Inches(0.82),
            top=card_top + Inches(0.38),
            width=Inches(1.64),
            height=Inches(0.78),
            text=str(event.get("detail") or event.get("title") or "")[:120],
            font_pt=max(9, palette.body_pt - 3),
            rgb=palette.body_rgb,
        )
        stem = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            center_x,
            lane_mid_y,
            center_x,
            card_top + (Inches(1.32) if idx % 2 == 0 else 0),
        )
        stem.line.color.rgb = RGBColor(*palette.card_border_rgb)
        stem.line.width = Inches(0.015)


def _add_status_board(
    slide: Any,
    *,
    vis: dict[str, Any],
    bullets: list[tuple[str, int]],
    palette: PptxSlidePalette,
) -> None:
    from pptx.util import Inches  # type: ignore[import-untyped]

    cards_raw = vis.get("cards") if isinstance(vis.get("cards"), list) else []
    cards = [c for c in cards_raw[:4] if isinstance(c, dict)]
    if not cards:
        _add_bullet_box(
            slide,
            left=Inches(0.72),
            top=Inches(1.88),
            width=Inches(11.9),
            height=Inches(4.6),
            items=bullets or [("Add status cards to pptx_spec.visual.cards.", 0)],
            palette=palette,
        )
        return

    if bullets:
        _add_bullet_box(
            slide,
            left=Inches(0.72),
            top=Inches(1.88),
            width=Inches(3.15),
            height=Inches(4.55),
            items=bullets[:6],
            palette=palette,
        )
        board_left = Inches(4.12)
        board_width = Inches(8.5)
    else:
        board_left = Inches(0.72)
        board_width = Inches(11.9)
    gap = Inches(0.22)
    cols = 2 if len(cards) > 2 else len(cards)
    rows = 2 if len(cards) > 2 else 1
    card_w = (board_width - gap * (cols - 1)) / max(cols, 1)
    card_h = Inches(2.05 if rows == 2 else 4.1)
    for idx, card in enumerate(cards):
        row = idx // cols
        col = idx % cols
        left = board_left + col * (card_w + gap)
        top = Inches(1.88) + row * (card_h + gap)
        _add_card(slide, left=left, top=top, width=card_w, height=card_h, palette=palette)
        title = str(card.get("title") or f"Card {idx + 1}")[:50]
        status = str(card.get("status") or "").strip()
        _add_textbox(
            slide,
            left=left + Inches(0.16),
            top=top + Inches(0.14),
            width=card_w - Inches(0.32),
            height=Inches(0.24),
            text=title,
            font_pt=palette.body_pt,
            rgb=palette.title_rgb,
            bold=True,
        )
        if status:
            _add_textbox(
                slide,
                left=left + Inches(0.16),
                top=top + Inches(0.42),
                width=card_w - Inches(0.32),
                height=Inches(0.22),
                text=status[:60],
                font_pt=max(10, palette.body_pt - 2),
                rgb=palette.accent_rgb,
                bold=True,
            )
        _add_textbox(
            slide,
            left=left + Inches(0.16),
            top=top + Inches(0.72),
            width=card_w - Inches(0.32),
            height=card_h - Inches(0.88),
            text=str(card.get("detail") or card.get("summary") or "")[:180],
            font_pt=max(10, palette.body_pt - 2),
            rgb=palette.body_rgb,
        )


def _add_preview_comparison_side_panel(
    slide: Any,
    *,
    panel: PreviewComparisonPanel,
    palette: PptxSlidePalette,
) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore[import-untyped]
    from pptx.util import Inches, Pt  # type: ignore[import-untyped]

    card_left = Inches(8.55)
    card_width = Inches(4.05)

    _add_card(slide, left=card_left, top=Inches(1.88), width=card_width, height=Inches(1.86), palette=palette)
    _add_textbox(
        slide,
        left=card_left + Inches(0.18),
        top=Inches(2.02),
        width=card_width - Inches(0.36),
        height=Inches(0.24),
        text="能力覆盖度",
        font_pt=palette.body_pt + 1,
        rgb=palette.title_rgb,
        bold=True,
    )
    for idx, (label, value) in enumerate(panel.coverage_rows[:3]):
        top = Inches(2.42 + idx * 0.46)
        _add_textbox(
            slide,
            left=card_left + Inches(0.18),
            top=top,
            width=Inches(2.35),
            height=Inches(0.16),
            text=label[:60],
            font_pt=max(10, palette.body_pt - 1),
            rgb=palette.body_rgb,
        )
        _add_textbox(
            slide,
            left=card_left + Inches(2.72),
            top=top,
            width=Inches(0.95),
            height=Inches(0.16),
            text=value[:20],
            font_pt=max(10, palette.body_pt - 1),
            rgb=palette.body_rgb,
            bold=True,
        )
        rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, card_left + Inches(0.18), top + Inches(0.23), Inches(3.45), Inches(0.12))
        rail.adjustments[0] = 0.2
        rail.fill.solid()
        rail.fill.fore_color.rgb = RGBColor(226, 232, 240)
        rail.line.color.rgb = RGBColor(226, 232, 240)
        bar_width = Inches(3.45)
        pct = 1.0 if "100" in value else 0.3
        fill = palette.success_rgb if pct >= 1.0 else (148, 163, 184)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, card_left + Inches(0.18), top + Inches(0.23), bar_width * pct, Inches(0.12))
        bar.adjustments[0] = 0.2
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*fill)
        bar.line.color.rgb = RGBColor(*fill)

    _add_card(slide, left=card_left, top=Inches(3.95), width=card_width, height=Inches(2.02), palette=palette)
    _add_textbox(
        slide,
        left=card_left + Inches(0.18),
        top=Inches(4.08),
        width=card_width - Inches(0.36),
        height=Inches(0.24),
        text="独家能力示意",
        font_pt=palette.body_pt + 1,
        rgb=palette.title_rgb,
        bold=True,
    )
    for idx, item in enumerate(panel.insight_rows[:4]):
        top = Inches(4.46 + idx * 0.42)
        dot_color = palette.success_rgb if item.startswith("OT") else palette.danger_rgb
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, card_left + Inches(0.18), top + Inches(0.03), Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(*dot_color)
        dot.line.color.rgb = RGBColor(*dot_color)
        box = slide.shapes.add_textbox(card_left + Inches(0.34), top, Inches(3.45), Inches(0.22))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        if "：" in item:
            left_label, right_text = item.split("：", 1)
            r1 = p.add_run()
            r1.text = left_label + "："
            r1.font.bold = True
            r1.font.size = Pt(max(10, palette.body_pt - 1))
            r1.font.color.rgb = RGBColor(*palette.body_rgb)
            r2 = p.add_run()
            r2.text = right_text.strip()[:120]
            r2.font.size = Pt(max(10, palette.body_pt - 1))
            r2.font.color.rgb = RGBColor(*palette.body_rgb)
        else:
            p.text = item[:140]
            p.font.size = Pt(max(10, palette.body_pt - 1))
            p.font.color.rgb = RGBColor(*palette.body_rgb)


def _add_takeaway_banner(slide: Any, *, text: str, palette: PptxSlidePalette) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore[import-untyped]
    from pptx.util import Inches, Pt  # type: ignore[import-untyped]

    if not text.strip():
        return
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(6.22), Inches(11.9), Inches(0.54))
    shape.adjustments[0] = 0.18
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 253, 244)
    shape.line.color.rgb = RGBColor(187, 247, 208)
    shape.line.width = Inches(0.01)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text[:220]
    p.font.size = Pt(max(11, palette.body_pt))
    p.font.bold = True
    p.font.color.rgb = RGBColor(22, 101, 52)


def _add_footer_line(slide: Any, text: str, palette: PptxSlidePalette) -> None:
    from pptx.util import Inches  # type: ignore[import-untyped]

    if not text.strip():
        return
    _add_textbox(
        slide,
        left=Inches(0.68),
        top=Inches(6.88),
        width=Inches(11.6),
        height=Inches(0.24),
        text=text[:320],
        font_pt=max(9, palette.body_pt - 3),
        rgb=palette.muted_rgb,
    )


def build_pptx_with_specs(session_id: str, framework: dict[str, Any], dest: Path) -> None:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.util import Inches  # type: ignore[import-untyped]
    except ImportError as e:  # pragma: no cover
        raise RuntimeError("python-pptx is required for PPTX export. pip install python-pptx") from e

    dest.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    fw_style = _normalize_deck_style(str(framework.get("deck_style")))
    slides = _slides_list(framework)
    if not slides:
        slides = [{"id": "__empty__", "title": "Empty deck", "subtitle": "", "main": "Add slides to the framework."}]
    blank_layout = prs.slide_layouts[6]
    surface = _surface_variant(session_id)

    with temp_media_dir() as tmp_name:
        tmp_dir = Path(tmp_name)
        for fw_slide in slides:
            sid = str(fw_slide.get("id") or "")
            meta = _load_meta(session_id, sid)
            scene = meta.get("pptx_scene") if isinstance(meta, dict) else None
            slide = prs.slides.add_slide(blank_layout)
            if has_scene(scene):
                render_scene_to_pptx_slide(slide, scene, tmp_dir)
                continue
            _render_spec_slide(slide, session_id, sid, fw_slide, meta, fw_style, surface)

    prs.save(str(dest))


def _render_spec_slide(
    slide: Any,
    session_id: str,
    sid: str,
    fw_slide: dict[str, Any],
    meta: dict[str, Any],
    fw_style: str,
    surface: str,
) -> None:
    from pptx.util import Inches  # type: ignore[import-untyped]

    preview_comparison_panel = _extract_preview_comparison_panel(session_id, sid)
    raw_spec = meta.get("pptx_spec") if meta else None
    spec = _coerce_spec(raw_spec)
    view = _merged_slide_view(fw_slide, spec, fw_style)
    palette = _palette_for(view["deck_style"], surface)
    _apply_background(slide, palette.slide_bg)

    layout_name = view["layout"]
    title_font = _fit_font_size(
        view["title"],
        palette.title_pt,
        thresholds=[(60, palette.title_pt - 2), (90, palette.title_pt - 4), (120, palette.title_pt - 6)],
    )
    _add_textbox(
        slide,
        left=Inches(0.72),
        top=Inches(0.5),
        width=Inches(11.8),
        height=Inches(0.75),
        text=view["title"][:500],
        font_pt=title_font,
        rgb=palette.title_rgb,
        bold=True,
    )
    if view["subtitle"]:
        _add_textbox(
            slide,
            left=Inches(0.74),
            top=Inches(1.28),
            width=Inches(11.4),
            height=Inches(0.38),
            text=view["subtitle"][:260],
            font_pt=palette.subtitle_pt,
            rgb=palette.subtitle_rgb,
        )

    if layout_name == "title_only":
        _add_footer_line(slide, view["footer_source"], palette)
        return

    if layout_name == "two_column":
        _add_bullet_box(
            slide,
            left=Inches(0.72),
            top=Inches(1.85),
            width=Inches(5.8),
            height=Inches(4.55),
            items=view["bullets"],
            palette=palette,
        )
        rb: list[tuple[str, int]] = list(view["right_col"])
        if not rb and view["takeaway"]:
            rb = [(view["takeaway"], 0)]
        _add_bullet_box(
            slide,
            left=Inches(6.82),
            top=Inches(1.85),
            width=Inches(5.75),
            height=Inches(4.55),
            items=rb or [(" ", 0)],
            palette=palette,
        )
        _add_footer_line(slide, view["footer_source"], palette)
        return

    vk = str((view.get("visual") or {}).get("kind") or "none").lower()
    vis = view.get("visual") if isinstance(view.get("visual"), dict) else {}
    if layout_name == "comparison_board" or vk == "comparison_cards":
        _add_comparison_board(slide, vis=vis, bullets=view["bullets"], palette=palette)
    elif layout_name == "process_flow" or vk == "process_flow":
        _add_process_flow(slide, vis=vis, bullets=view["bullets"], palette=palette)
    elif layout_name == "matrix_2x2" or vk == "matrix":
        _add_matrix_2x2(slide, vis=vis, bullets=view["bullets"], palette=palette)
    elif layout_name == "timeline" or vk == "timeline":
        _add_timeline(slide, vis=vis, bullets=view["bullets"], palette=palette)
    elif layout_name == "status_board" or vk == "status_cards":
        _add_status_board(slide, vis=vis, bullets=view["bullets"], palette=palette)
    elif layout_name == "content_table" or vk == "table":
        rowdata: list[list[str]] = []
        hraw = vis.get("headers")
        hdr = [str(x) for x in hraw][:12] if isinstance(hraw, list) else []
        if not hdr:
            hdr = ["Column A", "Column B"]
        raw_rows = vis.get("rows") if isinstance(vis.get("rows"), list) else []
        for row in raw_rows[:16]:
            if isinstance(row, list):
                rowdata.append([str(c)[:200] for c in row])
        if preview_comparison_panel is not None:
            _add_card(
                slide,
                left=Inches(0.72),
                top=Inches(1.88),
                width=Inches(7.58),
                height=Inches(4.1),
                palette=palette,
            )
            _add_table_to_slide(
                slide,
                hdr,
                rowdata,
                palette,
                left=Inches(0.94),
                top=Inches(2.08),
                width=Inches(7.14),
                height=Inches(3.72),
                highlight_ot=True,
            )
            _add_preview_comparison_side_panel(slide, panel=preview_comparison_panel, palette=palette)
        elif bool(view["bullets"]):
            _add_bullet_box(
                slide,
                left=Inches(0.72),
                top=Inches(1.88),
                width=Inches(4.3),
                height=Inches(4.45),
                items=view["bullets"],
                palette=palette,
            )
            _add_card(
                slide,
                left=Inches(5.22),
                top=Inches(1.88),
                width=Inches(7.4),
                height=Inches(4.45),
                palette=palette,
            )
            _add_table_to_slide(
                slide,
                hdr,
                rowdata,
                palette,
                left=Inches(5.45),
                top=Inches(2.12),
                width=Inches(6.96),
                height=Inches(3.9),
            )
        else:
            _add_card(
                slide,
                left=Inches(0.72),
                top=Inches(1.88),
                width=Inches(11.9),
                height=Inches(4.45),
                palette=palette,
            )
            _add_table_to_slide(
                slide,
                hdr,
                rowdata,
                palette,
                left=Inches(0.95),
                top=Inches(2.12),
                width=Inches(11.46),
                height=Inches(3.9),
            )
    elif layout_name == "kpi_strip" and isinstance(view["metrics"], list) and view["metrics"]:
        metrics = [m for m in view["metrics"][:4] if isinstance(m, dict)]
        _add_bullet_box(
            slide,
            left=Inches(0.72),
            top=Inches(2.75),
            width=Inches(11.9),
            height=Inches(2.6),
            items=view["bullets"],
            palette=palette,
        )
        for idx, m in enumerate(metrics):
            left = Inches(0.72 + idx * 3.0)
            _add_card(slide, left=left, top=Inches(1.82), width=Inches(2.7), height=Inches(0.76), palette=palette)
            _add_textbox(
                slide,
                left=left + Inches(0.16),
                top=Inches(1.95),
                width=Inches(2.38),
                height=Inches(0.22),
                text=str(m.get("label") or "")[:60],
                font_pt=max(10, palette.body_pt - 2),
                rgb=palette.subtitle_rgb,
                bold=True,
            )
            _add_textbox(
                slide,
                left=left + Inches(0.16),
                top=Inches(2.18),
                width=Inches(2.38),
                height=Inches(0.26),
                text=str(m.get("value") or "")[:40],
                font_pt=palette.body_pt + 2,
                rgb=palette.accent_rgb,
                bold=True,
            )
    else:
        body_height = Inches(4.5 if view["takeaway"] else 5.0)
        _add_bullet_box(
            slide,
            left=Inches(0.72),
            top=Inches(1.88),
            width=Inches(11.9),
            height=body_height,
            items=view["bullets"] or [(view["subtitle"], 0)] if view["subtitle"] else [(" ", 0)],
            palette=palette,
        )

    if preview_comparison_panel is not None and (preview_comparison_panel.takeaway or view["takeaway"]):
        _add_takeaway_banner(slide, text=preview_comparison_panel.takeaway or view["takeaway"], palette=palette)
    elif view["takeaway"]:
        _add_textbox(
            slide,
            left=Inches(0.78),
            top=Inches(6.4),
            width=Inches(11.4),
            height=Inches(0.28),
            text=f"► {view['takeaway'][:220]}",
            font_pt=palette.body_pt,
            rgb=palette.accent_rgb,
            bold=True,
        )
    _add_footer_line(slide, view["footer_source"], palette)
