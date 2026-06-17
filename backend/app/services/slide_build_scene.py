"""
Element-level slide scene rendering for editable PPTX export.

The scene uses a fixed 1600x900 canvas. Coordinates are converted directly to
PowerPoint inches, which keeps browser preview and PPTX export on the same grid.
"""

from __future__ import annotations

import base64
import re
import tempfile
from pathlib import Path
from typing import Any
from urllib.parse import unquote, urlparse
from xml.etree import ElementTree

CANVAS_W = 1600.0
CANVAS_H = 900.0
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


def has_scene(scene: Any) -> bool:
    return isinstance(scene, dict) and isinstance(scene.get("elements"), list)


def render_scene_to_pptx_slide(slide: Any, scene: dict[str, Any], tmp_dir: Path) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.util import Inches  # type: ignore[import-untyped]

    bg = _rgb(scene.get("background") or scene.get("bg") or "#ffffff")
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg)

    elements = scene.get("elements")
    if not isinstance(elements, list):
        return
    for idx, raw in enumerate(elements):
        if not isinstance(raw, dict):
            continue
        etype = str(raw.get("type") or "").strip().lower()
        try:
            if etype in {"rect", "round_rect", "rounded_rect", "card"}:
                _add_rect(slide, raw, rounded=etype != "rect")
            elif etype in {"ellipse", "circle", "oval"}:
                _add_ellipse(slide, raw)
            elif etype in {"line", "connector"}:
                _add_line(slide, raw)
            elif etype == "text":
                _add_text(slide, raw)
            elif etype == "table":
                _add_table(slide, raw)
            elif etype == "image":
                _add_image(slide, raw, tmp_dir, idx)
            elif etype == "svg":
                _add_svg_primitives(slide, raw, tmp_dir, idx)
        except Exception:
            # Keep export resilient: one unsupported scene object should not
            # prevent the rest of the slide from staying editable.
            continue


def _in(value: Any, axis: str = "x") -> Any:
    from pptx.util import Inches  # type: ignore[import-untyped]

    scale = SLIDE_W_IN / CANVAS_W if axis == "x" else SLIDE_H_IN / CANVAS_H
    try:
        v = float(value)
    except (TypeError, ValueError):
        v = 0.0
    return Inches(v * scale)


def _rgb(value: Any, fallback: tuple[int, int, int] = (255, 255, 255)) -> tuple[int, int, int]:
    if isinstance(value, (list, tuple)) and len(value) >= 3:
        try:
            return (int(value[0]), int(value[1]), int(value[2]))
        except (TypeError, ValueError):
            return fallback
    s = str(value or "").strip()
    if not s or s.lower() in {"none", "transparent"}:
        return fallback
    if s.startswith("#"):
        h = s[1:]
        if len(h) == 3:
            h = "".join(ch * 2 for ch in h)
        if len(h) >= 6:
            try:
                return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            except ValueError:
                return fallback
    m = re.match(r"rgba?\(([^)]+)\)", s, flags=re.I)
    if m:
        parts = [p.strip() for p in m.group(1).split(",")]
        if len(parts) >= 3:
            try:
                return (int(float(parts[0])), int(float(parts[1])), int(float(parts[2])))
            except ValueError:
                return fallback
    return fallback


def _num(value: Any, default: float = 0.0) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        m = re.search(r"-?\d+(?:\.\d+)?", str(value or ""))
        if not m:
            return default
        try:
            return float(m.group(0))
        except ValueError:
            return default


def _bounds(el: dict[str, Any]) -> tuple[Any, Any, Any, Any]:
    return (
        _in(el.get("x", el.get("left", 0)), "x"),
        _in(el.get("y", el.get("top", 0)), "y"),
        _in(el.get("w", el.get("width", 0)), "x"),
        _in(el.get("h", el.get("height", 0)), "y"),
    )


def _apply_fill_and_line(shape: Any, el: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.util import Pt  # type: ignore[import-untyped]

    fill_color = str(el.get("fill", "")).strip()
    if fill_color.lower() in {"", "none", "transparent"}:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*_rgb(fill_color))

    stroke = str(el.get("stroke", el.get("line", ""))).strip()
    if stroke.lower() in {"", "none", "transparent"}:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = RGBColor(*_rgb(stroke, (0, 0, 0)))
        shape.line.width = Pt(max(0.25, _num(el.get("strokeWidth", el.get("lineWidth", 1)), 1)))


def _add_rect(slide: Any, el: dict[str, Any], *, rounded: bool) -> None:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore[import-untyped]

    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, *_bounds(el))
    if rounded:
        radius = max(0.0, min(0.5, _num(el.get("radius", el.get("rx", 18)), 18) / 100.0))
        shape.adjustments[0] = radius
    _apply_fill_and_line(shape, el)


def _add_ellipse(slide: Any, el: dict[str, Any]) -> None:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore[import-untyped]

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, *_bounds(el))
    _apply_fill_and_line(shape, el)


def _add_line(slide: Any, el: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.shapes import MSO_CONNECTOR  # type: ignore[import-untyped]
    from pptx.util import Pt  # type: ignore[import-untyped]

    x1 = _in(el.get("x1", el.get("x", 0)), "x")
    y1 = _in(el.get("y1", el.get("y", 0)), "y")
    x2 = _in(el.get("x2", _num(el.get("x", 0)) + _num(el.get("w", 0))), "x")
    y2 = _in(el.get("y2", _num(el.get("y", 0)) + _num(el.get("h", 0))), "y")
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = RGBColor(*_rgb(el.get("stroke") or el.get("color") or "#000000", (0, 0, 0)))
    line.line.width = Pt(max(0.25, _num(el.get("strokeWidth", el.get("width", 1)), 1)))


def _add_text(slide: Any, el: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # type: ignore[import-untyped]
    from pptx.util import Pt  # type: ignore[import-untyped]

    box = slide.shapes.add_textbox(*_bounds(el))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = bool(el.get("wrap", True))
    valign = str(el.get("valign") or el.get("verticalAlign") or "top").lower()
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    text = str(el.get("text") or "")
    lines = text.splitlines() or [""]
    color = RGBColor(*_rgb(el.get("color") or "#111827", (17, 24, 39)))
    size = Pt(max(1, _num(el.get("fontSize", el.get("size", 18)), 18)))
    align = str(el.get("align") or "left").lower()
    align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.bold = bool(el.get("bold"))
        p.font.italic = bool(el.get("italic"))
        p.font.color.rgb = color
        if align in align_map:
            p.alignment = align_map[align]


def _add_table(slide: Any, el: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.util import Pt  # type: ignore[import-untyped]

    headers = el.get("headers") if isinstance(el.get("headers"), list) else []
    rows = el.get("rows") if isinstance(el.get("rows"), list) else []
    body = [r for r in rows if isinstance(r, list)]
    cols = max(len(headers), *(len(r) for r in body), 1)
    nrows = max(1, len(body) + (1 if headers else 0))
    graphic = slide.shapes.add_table(nrows, cols, *_bounds(el))
    tbl = graphic.table
    font_size = Pt(max(6, _num(el.get("fontSize", 13), 13)))
    body_rgb = RGBColor(*_rgb(el.get("color") or "#334155", (51, 65, 85)))
    header_fill = RGBColor(*_rgb(el.get("headerFill") or "#e2e8f0", (226, 232, 240)))
    header_rgb = RGBColor(*_rgb(el.get("headerColor") or "#0f172a", (15, 23, 42)))
    start = 0
    if headers:
        start = 1
        for c in range(cols):
            cell = tbl.cell(0, c)
            cell.text = str(headers[c] if c < len(headers) else "")[:500]
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = font_size
            p.font.color.rgb = header_rgb
    for ridx, row in enumerate(body[: nrows - start], start=start):
        for c in range(cols):
            cell = tbl.cell(ridx, c)
            cell.text = str(row[c] if c < len(row) else "")[:1000]
            p = cell.text_frame.paragraphs[0]
            p.font.size = font_size
            p.font.color.rgb = body_rgb


def _resolve_media_path(src: Any, tmp_dir: Path, idx: int) -> Path | None:
    s = str(src or "").strip()
    if not s:
        return None
    if s.startswith("data:image/"):
        head, _, payload = s.partition(",")
        ext = "png"
        m = re.search(r"data:image/([a-zA-Z0-9+.-]+);base64", head)
        if m:
            ext = "jpg" if m.group(1).lower() == "jpeg" else m.group(1).lower()
        try:
            data = base64.b64decode(payload)
        except Exception:
            return None
        out = tmp_dir / f"scene_media_{idx}.{ext}"
        out.write_bytes(data)
        return out
    if s.startswith("file://"):
        return Path(unquote(urlparse(s).path))
    p = Path(s)
    return p if p.is_file() else None


def _add_image(slide: Any, el: dict[str, Any], tmp_dir: Path, idx: int) -> None:
    path = _resolve_media_path(el.get("src") or el.get("path"), tmp_dir, idx)
    if not path or not path.is_file():
        return
    left, top, width, height = _bounds(el)
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def _add_svg_primitives(slide: Any, el: dict[str, Any], tmp_dir: Path, idx: int) -> None:
    svg = str(el.get("svg") or "").strip()
    if not svg:
        _add_image(slide, el, tmp_dir, idx)
        return
    try:
        root = ElementTree.fromstring(svg)
    except ElementTree.ParseError:
        return
    base_x = _num(el.get("x", 0))
    base_y = _num(el.get("y", 0))
    vb = [float(x) for x in re.findall(r"-?\d+(?:\.\d+)?", root.get("viewBox") or "")]
    src_w = vb[2] if len(vb) == 4 and vb[2] else _num(root.get("width"), CANVAS_W)
    src_h = vb[3] if len(vb) == 4 and vb[3] else _num(root.get("height"), CANVAS_H)
    sx = _num(el.get("w", CANVAS_W), CANVAS_W) / max(src_w, 1)
    sy = _num(el.get("h", CANVAS_H), CANVAS_H) / max(src_h, 1)
    for node in root.iter():
        tag = node.tag.rsplit("}", 1)[-1].lower()
        style = _style_attrs(node)
        if tag == "rect":
            _add_rect(slide, {
                "type": "round_rect" if _num(node.get("rx"), 0) or _num(node.get("ry"), 0) else "rect",
                "x": base_x + _num(node.get("x"), 0) * sx,
                "y": base_y + _num(node.get("y"), 0) * sy,
                "w": _num(node.get("width"), 0) * sx,
                "h": _num(node.get("height"), 0) * sy,
                "rx": _num(node.get("rx"), 0) * sx,
                **style,
            }, rounded=bool(_num(node.get("rx"), 0) or _num(node.get("ry"), 0)))
        elif tag in {"circle", "ellipse"}:
            cx = _num(node.get("cx"), 0) * sx
            cy = _num(node.get("cy"), 0) * sy
            rx = (_num(node.get("r"), 0) or _num(node.get("rx"), 0)) * sx
            ry = (_num(node.get("r"), 0) or _num(node.get("ry"), 0)) * sy
            _add_ellipse(slide, {"x": base_x + cx - rx, "y": base_y + cy - ry, "w": rx * 2, "h": ry * 2, **style})
        elif tag == "line":
            _add_line(slide, {
                "x1": base_x + _num(node.get("x1"), 0) * sx,
                "y1": base_y + _num(node.get("y1"), 0) * sy,
                "x2": base_x + _num(node.get("x2"), 0) * sx,
                "y2": base_y + _num(node.get("y2"), 0) * sy,
                **style,
            })
        elif tag == "text":
            _add_text(slide, {
                "x": base_x + _num(node.get("x"), 0) * sx,
                "y": base_y + _num(node.get("y"), 0) * sy,
                "w": _num(node.get("width"), 240) * sx,
                "h": _num(node.get("height"), 48) * sy,
                "text": "".join(node.itertext()).strip(),
                "fontSize": _num(style.get("fontSize"), 18),
                "color": style.get("fill") or style.get("color") or "#111827",
            })
        elif tag == "path":
            _add_simple_path(slide, node.get("d") or "", base_x, base_y, sx, sy, style)


def _style_attrs(node: ElementTree.Element) -> dict[str, Any]:
    out: dict[str, Any] = {}
    for key in ("fill", "stroke", "stroke-width", "color", "font-size"):
        if node.get(key) is not None:
            clean = key.replace("stroke-width", "strokeWidth").replace("font-size", "fontSize")
            out[clean] = node.get(key)
    style = node.get("style") or ""
    for part in style.split(";"):
        k, sep, v = part.partition(":")
        if not sep:
            continue
        clean = k.strip().replace("stroke-width", "strokeWidth").replace("font-size", "fontSize")
        out[clean] = v.strip()
    return out


def _add_simple_path(slide: Any, d: str, base_x: float, base_y: float, sx: float, sy: float, style: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.util import Pt  # type: ignore[import-untyped]

    tokens = re.findall(r"[MLZmlz]|-?\d+(?:\.\d+)?", d)
    if not tokens:
        return
    points: list[tuple[float, float]] = []
    i = 0
    cmd = ""
    while i < len(tokens):
        tok = tokens[i]
        if re.match(r"[MLZmlz]", tok):
            cmd = tok
            i += 1
            if cmd.lower() == "z":
                break
            continue
        if cmd.lower() not in {"m", "l"} or i + 1 >= len(tokens):
            return
        x = base_x + float(tokens[i]) * sx
        y = base_y + float(tokens[i + 1]) * sy
        points.append((x, y))
        i += 2
    if len(points) < 2:
        return
    ff = slide.shapes.build_freeform(_in(points[0][0], "x"), _in(points[0][1], "y"))
    segs = [(_in(x, "x"), _in(y, "y")) for x, y in points[1:]]
    ff.add_line_segments(segs, close=bool("z" in d.lower()))
    shape = ff.convert_to_shape()
    fill = str(style.get("fill", "none")).strip().lower()
    if fill not in {"", "none", "transparent"}:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*_rgb(fill))
    else:
        shape.fill.background()
    stroke = str(style.get("stroke", "#000000")).strip()
    if stroke.lower() not in {"", "none", "transparent"}:
        shape.line.color.rgb = RGBColor(*_rgb(stroke, (0, 0, 0)))
        shape.line.width = Pt(max(0.25, _num(style.get("strokeWidth"), 1)))


def temp_media_dir() -> tempfile.TemporaryDirectory[str]:
    return tempfile.TemporaryDirectory(prefix="unbox-pptx-scene-")
