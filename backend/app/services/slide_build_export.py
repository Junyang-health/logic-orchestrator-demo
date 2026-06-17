"""
Build PPTX / PDF deck files.

When generated HTML previews exist, PDF export is rendered from those previews so
the downloaded PDF matches what the user saw in the browser. PPTX export defaults
to the native structured renderer so the deck remains editable.
"""

from __future__ import annotations

import os
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any

from app.services.slide_build_artifacts import (
    extract_slide_inner_from_document,
    export_pdf_path,
    export_pptx_path,
    read_slide_document,
    render_slide_document,
    slide_html_path,
)
from app.services.slide_export_review import export_review_path, review_export_consistency
from app.services.slide_build_pptx_render import build_pptx_with_specs, merged_slide_views

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
RENDER_W_PX = 1600
RENDER_H_PX = 900
CHROME_RENDER_TIMEOUT_SEC = 12


def _pdf_palette(deck_style: str) -> dict[str, tuple[float, float, float]]:
    s = (deck_style or "consulting_mbb").strip().lower()
    if s == "creative":
        return {
            "bg": (0.06, 0.09, 0.16),
            "title": (0.97, 0.98, 1.0),
            "sub": (0.8, 0.84, 0.92),
            "body": (0.88, 0.91, 0.96),
            "muted": (0.58, 0.64, 0.72),
        }
    if s == "academic":
        return {
            "bg": (0.99, 0.99, 0.98),
            "title": (0.13, 0.15, 0.17),
            "sub": (0.28, 0.29, 0.32),
            "body": (0.21, 0.21, 0.21),
            "muted": (0.44, 0.46, 0.5),
        }
    if s == "government":
        return {
            "bg": (1.0, 1.0, 1.0),
            "title": (0.09, 0.17, 0.3),
            "sub": (0.26, 0.32, 0.41),
            "body": (0.2, 0.24, 0.28),
            "muted": (0.47, 0.51, 0.56),
        }
    return {
        "bg": (1.0, 1.0, 1.0),
        "title": (0.12, 0.16, 0.22),
        "sub": (0.28, 0.33, 0.41),
        "body": (0.2, 0.25, 0.33),
        "muted": (0.5, 0.55, 0.62),
    }


def build_pptx_file(
    session_id: str,
    framework: dict[str, Any],
    dest: Path,
    rendered_previews: list[Path] | None = None,
) -> None:
    mode = (os.getenv("UNBOX_PPTX_EXPORT_MODE") or "editable").strip().lower()
    if mode in {"image", "preview", "raster"}:
        rendered = (
            _render_preview_pngs(session_id, framework)
            if rendered_previews is None
            else rendered_previews
        )
        if rendered:
            build_pptx_from_png_paths(rendered, dest)
            return
    build_pptx_with_specs(session_id, framework, dest)


def build_pdf_file(
    session_id: str,
    framework: dict[str, Any],
    dest: Path,
    rendered_previews: list[Path] | None = None,
) -> None:
    rendered = (
        _render_preview_pngs(session_id, framework)
        if rendered_previews is None
        else rendered_previews
    )
    if rendered:
        _build_pdf_from_pngs(rendered, dest)
        return
    _build_pdf_from_framework(session_id, framework, dest)


def _build_pdf_from_framework(session_id: str, framework: dict[str, Any], dest: Path) -> None:
    try:
        from reportlab.lib.pagesizes import landscape  # type: ignore[import-untyped]
        from reportlab.lib.units import inch  # type: ignore[import-untyped]
        from reportlab.pdfgen import canvas  # type: ignore[import-untyped]
    except ImportError as e:  # pragma: no cover
        raise RuntimeError("reportlab is required for PDF export. pip install reportlab") from e

    dest.parent.mkdir(parents=True, exist_ok=True)
    w, h = landscape((13.333 * inch, 7.5 * inch))
    c = canvas.Canvas(str(dest), pagesize=(w, h))

    views = merged_slide_views(session_id, framework)
    if not views:
        views = [{"title": "Empty deck", "subtitle": "", "bullets": [], "takeaway": "", "footer_source": "", "deck_style": "consulting_mbb"}]

    for view in views:
        palette = _pdf_palette(str(view.get("deck_style") or "consulting_mbb"))
        title = str(view.get("title") or "Slide").strip() or "Slide"
        subtitle = str(view.get("subtitle") or "").strip()
        takeaway = str(view.get("takeaway") or "").strip()
        footer = str(view.get("footer_source") or "").strip()

        c.setFillColorRGB(*palette["bg"])
        c.rect(0, 0, w, h, fill=1, stroke=0)

        c.setFillColorRGB(*palette["title"])
        c.setFont("Helvetica-Bold", 22)
        c.drawString(0.55 * inch, h - 0.9 * inch, title[:120])

        y = h - 1.35 * inch
        c.setFont("Helvetica", 12)
        if subtitle:
            c.setFillColorRGB(*palette["sub"])
            c.drawString(0.55 * inch, y, subtitle[:220])
            y -= 0.32 * inch

        c.setFillColorRGB(*palette["body"])
        c.setFont("Helvetica", 11)
        bl = view.get("bullets") if isinstance(view.get("bullets"), list) else []
        for item in bl:
            line = ""
            if isinstance(item, (list, tuple)) and len(item) >= 1:
                line = str(item[0]).strip()
            elif isinstance(item, dict):
                line = str(item.get("text") or "").strip()
            elif isinstance(item, str):
                line = item.strip()
            if not line:
                continue
            rest = line
            while rest:
                chunk, rest = rest[:98], rest[98:]
                c.drawString(0.55 * inch, y, chunk)
                y -= 0.2 * inch
                if y < 0.75 * inch:
                    break
            if y < 0.75 * inch:
                break

        if takeaway:
            c.setFillColorRGB(*palette["muted"])
            c.setFont("Helvetica-Bold", 10)
            yt = takeaway[:400]
            c.drawString(0.55 * inch, max(0.55 * inch, y - 0.15 * inch), yt)
            y = max(0.55 * inch, y - 0.4 * inch)

        if footer:
            c.setFillColorRGB(*palette["muted"])
            c.setFont("Helvetica", 9)
            c.drawString(0.55 * inch, 0.45 * inch, footer[:260])

        c.showPage()

    c.save()


def _framework_slide_ids(framework: dict[str, Any]) -> list[str]:
    slides = framework.get("slides")
    if not isinstance(slides, list):
        return []
    out: list[str] = []
    for idx, slide in enumerate(slides, start=1):
        if isinstance(slide, dict):
            sid = str(slide.get("id") or "").strip()
            out.append(sid or f"slide_{idx}")
    return out


def _preview_html_paths(session_id: str, framework: dict[str, Any]) -> list[Path]:
    paths: list[Path] = []
    for sid in _framework_slide_ids(framework):
        p = slide_html_path(session_id, sid)
        if p.is_file():
            paths.append(p)
    return paths


def _normalized_preview_html_paths(session_id: str, framework: dict[str, Any], out_dir: Path) -> list[Path]:
    paths: list[Path] = []
    normalized_dir = out_dir / "normalized_html"
    normalized_dir.mkdir(parents=True, exist_ok=True)
    for idx, sid in enumerate(_framework_slide_ids(framework), start=1):
        full = read_slide_document(session_id, sid)
        if not full:
            continue
        inner = extract_slide_inner_from_document(full)
        doc = render_slide_document(inner) if inner else full
        p = normalized_dir / f"slide_{idx:03d}.html"
        p.write_text(doc, encoding="utf-8")
        paths.append(p)
    return paths


def _chrome_executable() -> str | None:
    env = (os.getenv("UNBOX_CHROME_PATH") or "").strip()
    candidates = [
        env,
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
        "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
        "google-chrome",
        "chromium",
        "chromium-browser",
    ]
    for cand in candidates:
        if not cand:
            continue
        if "/" in cand:
            if Path(cand).is_file():
                return cand
            continue
        try:
            subprocess.run(
                [cand, "--version"],
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=5,
            )
            return cand
        except Exception:
            continue
    return None


def _render_preview_pngs(session_id: str, framework: dict[str, Any]) -> list[Path]:
    out_dir = export_pdf_path(session_id).parent / "preview_renders"
    out_dir.mkdir(parents=True, exist_ok=True)
    raw_html_paths = _preview_html_paths(session_id, framework)
    html_paths = _normalized_preview_html_paths(session_id, framework, out_dir)
    if not html_paths:
        html_paths = raw_html_paths
    if not html_paths:
        return []

    rendered = _render_preview_pngs_with_playwright(html_paths, out_dir)
    if rendered:
        return rendered

    rendered = _render_preview_pngs_with_chrome_cli(html_paths, out_dir)
    if rendered:
        return rendered

    cached = _cached_preview_pngs(out_dir, expected_count=len(html_paths), source_html_paths=raw_html_paths)
    if cached:
        _log_preview_export_fallback("Using cached preview screenshots because fresh preview rendering failed.")
        return cached

    return []


def _cached_preview_pngs(out_dir: Path, *, expected_count: int, source_html_paths: list[Path]) -> list[Path]:
    cached = [out_dir / f"slide_{idx:03d}.png" for idx in range(1, expected_count + 1)]
    if not cached or any((not p.is_file() or p.stat().st_size <= 0) for p in cached):
        return []
    if len(source_html_paths) == expected_count:
        for png_path, html_path in zip(cached, source_html_paths):
            if html_path.is_file() and png_path.stat().st_mtime < html_path.stat().st_mtime:
                return []
    return cached


def _render_preview_pngs_with_playwright(html_paths: list[Path], out_dir: Path) -> list[Path]:
    try:
        from playwright.sync_api import sync_playwright  # type: ignore[import-untyped]
    except Exception as e:
        _log_preview_export_fallback(f"Playwright is unavailable ({type(e).__name__}: {e}).")
        return []

    rendered: list[Path] = []
    try:
        with sync_playwright() as p:
            try:
                browser = p.chromium.launch(channel="chrome", headless=True)
            except Exception:
                browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": RENDER_W_PX, "height": RENDER_H_PX}, device_scale_factor=1)
            for idx, html_path in enumerate(html_paths, start=1):
                png_path = out_dir / f"slide_{idx:03d}.png"
                page.goto(html_path.resolve().as_uri(), wait_until="networkidle", timeout=CHROME_RENDER_TIMEOUT_SEC * 1000)
                page.screenshot(path=str(png_path), full_page=False, timeout=CHROME_RENDER_TIMEOUT_SEC * 1000)
                if not png_path.is_file() or png_path.stat().st_size <= 0:
                    _log_preview_export_fallback(
                        f"Playwright did not produce preview slide {idx}; trying Chrome CLI fallback."
                    )
                    browser.close()
                    return []
                rendered.append(png_path)
            browser.close()
        return rendered
    except Exception as e:
        _log_preview_export_fallback(
            f"Playwright could not render preview slides ({type(e).__name__}: {e}); trying Chrome CLI fallback."
        )
        return []


def _render_preview_pngs_with_chrome_cli(html_paths: list[Path], out_dir: Path) -> list[Path]:
    chrome = _chrome_executable()
    if not chrome:
        _log_preview_export_fallback("Chrome/Chromium was not found; using structured export fallback.")
        return []

    rendered: list[Path] = []

    with tempfile.TemporaryDirectory(prefix="unbox-chrome-") as profile:
        for idx, html_path in enumerate(html_paths, start=1):
            png_path = out_dir / f"slide_{idx:03d}.png"
            cmd = [
                chrome,
                "--headless=new",
                "--disable-gpu",
                "--no-first-run",
                "--no-default-browser-check",
                "--hide-scrollbars",
                "--allow-file-access-from-files",
                "--disable-background-networking",
                "--disable-sync",
                "--disable-extensions",
                "--disable-features=Translate,MediaRouter",
                f"--user-data-dir={profile}",
                f"--window-size={RENDER_W_PX},{RENDER_H_PX}",
                f"--screenshot={png_path}",
                html_path.resolve().as_uri(),
            ]
            try:
                subprocess.run(
                    cmd,
                    check=True,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                    timeout=CHROME_RENDER_TIMEOUT_SEC,
                )
            except Exception as e:
                _log_preview_export_fallback(
                    f"Chrome could not render preview slide {idx} ({type(e).__name__}: {e}); "
                    "using structured export fallback."
                )
                return []
            if png_path.is_file() and png_path.stat().st_size > 0:
                rendered.append(png_path)
            else:
                _log_preview_export_fallback(
                    f"Chrome did not produce preview slide {idx}; using structured export fallback."
                )
                return []
    return rendered


def _log_preview_export_fallback(message: str) -> None:
    print(f"[slide-build] {message}", file=sys.stderr)


def build_pptx_from_png_paths(rendered: list[Path], dest: Path) -> None:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.util import Inches  # type: ignore[import-untyped]
    except ImportError as e:  # pragma: no cover
        raise RuntimeError("python-pptx is required for PPTX export. pip install python-pptx") from e

    dest.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    for png_path in rendered:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(dest)


def _build_pdf_from_pngs(rendered: list[Path], dest: Path) -> None:
    try:
        from reportlab.lib.units import inch  # type: ignore[import-untyped]
        from reportlab.lib.utils import ImageReader  # type: ignore[import-untyped]
        from reportlab.pdfgen import canvas  # type: ignore[import-untyped]
    except ImportError as e:  # pragma: no cover
        raise RuntimeError("reportlab is required for PDF export. pip install reportlab") from e

    dest.parent.mkdir(parents=True, exist_ok=True)
    w, h = SLIDE_W_IN * inch, SLIDE_H_IN * inch
    c = canvas.Canvas(str(dest), pagesize=(w, h))
    for png_path in rendered:
        c.drawImage(
            ImageReader(str(png_path)),
            0,
            0,
            width=w,
            height=h,
            preserveAspectRatio=False,
            mask="auto",
        )
        c.showPage()
    c.save()


def export_deck_files(session_id: str, framework: dict[str, Any]) -> dict[str, str]:
    """Write deck.pptx + deck.pdf; returns relative paths under data/."""
    from app.services.slide_build_artifacts import (
        DATA_DIR,
        ensure_session_dirs,
        export_pdf_path,
        export_pptx_path,
    )

    ensure_session_dirs(session_id)
    pptx_p = export_pptx_path(session_id)
    pdf_p = export_pdf_path(session_id)
    rendered = _render_preview_pngs(session_id, framework)
    build_pptx_file(session_id, framework, pptx_p)
    build_pdf_file(session_id, framework, pdf_p, rendered_previews=rendered)
    review = review_export_consistency(
        session_id,
        framework,
        preview_pngs=rendered,
        pptx_path=pptx_p,
        pdf_path=pdf_p,
    )

    def rel(p: Path) -> str:
        try:
            return str(p.relative_to(DATA_DIR))
        except ValueError:
            return str(p)

    review_p = export_review_path(session_id)
    return {
        "pptx": rel(pptx_p),
        "pdf": rel(pdf_p),
        "review": rel(review_p),
        "review_status": str(review.get("status") or ""),
    }
